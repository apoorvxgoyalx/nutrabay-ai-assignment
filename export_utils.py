"""
export_utils.py
---------------
Handles exporting the generated training module to PDF and PPTX formats.
"""

import io
from typing import Optional


# ── PDF Export ────────────────────────────────────────────────────────────────

def export_to_pdf(data: dict) -> bytes:
    """
    Export training module to a styled PDF using reportlab.

    Args:
        data: Parsed training module dict.

    Returns:
        PDF as bytes.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer,
            Table, TableStyle, HRFlowable
        )
        from reportlab.lib.enums import TA_LEFT, TA_CENTER
    except ImportError:
        raise ImportError("reportlab not installed. Run: pip install reportlab")

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=2 * cm,
        leftMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Title"],
        fontSize=22,
        textColor=colors.HexColor("#1a1a2e"),
        spaceAfter=6,
    )
    h1_style = ParagraphStyle(
        "CustomH1",
        parent=styles["Heading1"],
        fontSize=15,
        textColor=colors.HexColor("#16213e"),
        spaceBefore=14,
        spaceAfter=4,
        borderPad=4,
    )
    h2_style = ParagraphStyle(
        "CustomH2",
        parent=styles["Heading2"],
        fontSize=12,
        textColor=colors.HexColor("#0f3460"),
        spaceBefore=8,
        spaceAfter=3,
    )
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        spaceAfter=4,
    )
    bullet_style = ParagraphStyle(
        "CustomBullet",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        leftIndent=16,
        spaceAfter=3,
        bulletIndent=4,
    )

    story = []

    # Title
    title = data.get("document_title", "SOP Training Module")
    story.append(Paragraph(title, title_style))
    story.append(Paragraph("AI-Generated Training Module", styles["Italic"]))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#0f3460")))
    story.append(Spacer(1, 0.3 * cm))

    # ── Summary ──────────────────────────────────────────────────────────────
    summary = data.get("summary", {})
    story.append(Paragraph("📋 Summary", h1_style))
    story.append(Paragraph(summary.get("overview", ""), body_style))

    if summary.get("key_objectives"):
        story.append(Paragraph("Key Objectives", h2_style))
        for obj in summary["key_objectives"]:
            story.append(Paragraph(f"• {obj}", bullet_style))

    if summary.get("important_rules"):
        story.append(Paragraph("Important Rules", h2_style))
        for rule in summary["important_rules"]:
            story.append(Paragraph(f"• {rule}", bullet_style))

    story.append(Spacer(1, 0.4 * cm))

    # ── Training Steps ───────────────────────────────────────────────────────
    story.append(Paragraph("🎓 Training Steps", h1_style))
    for step in data.get("training_steps", []):
        num = step.get("step_number", "")
        step_title = step.get("title", "")
        story.append(Paragraph(f"Step {num}: {step_title}", h2_style))
        story.append(Paragraph(step.get("description", ""), body_style))
        if step.get("example"):
            story.append(
                Paragraph(f"<i>Example: {step['example']}</i>", body_style)
            )
        if step.get("tips"):
            for tip in step["tips"]:
                story.append(Paragraph(f"💡 {tip}", bullet_style))
        story.append(Spacer(1, 0.2 * cm))

    # ── Quiz ─────────────────────────────────────────────────────────────────
    story.append(Paragraph("📝 Quiz", h1_style))
    for q in data.get("quiz", []):
        qnum = q.get("question_number", "")
        story.append(Paragraph(f"Q{qnum}. {q.get('question', '')}", h2_style))
        for opt in q.get("options", []):
            story.append(Paragraph(opt, bullet_style))
        story.append(
            Paragraph(
                f"<b>Answer:</b> {q.get('answer', '')} — {q.get('explanation', '')}",
                body_style,
            )
        )
        story.append(Spacer(1, 0.2 * cm))

    doc.build(story)
    buffer.seek(0)
    return buffer.read()


# ── PPTX Export ───────────────────────────────────────────────────────────────

def export_to_pptx(data: dict) -> bytes:
    """
    Export training module to a styled PowerPoint presentation.

    Args:
        data: Parsed training module dict.

    Returns:
        PPTX as bytes.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    ACCENT = RGBColor(0x0F, 0x34, 0x60)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xE8, 0xF4, 0xFF)

    blank_layout = prs.slide_layouts[6]  # Blank

    def add_slide_with_header(title_text: str, color=DARK) -> object:
        slide = prs.slides.add_slide(blank_layout)
        # Header bar
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        tf = header.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        tf.word_wrap = False
        # Indent text
        p.space_before = Pt(6)
        return slide

    def add_text_box(slide, text, left, top, width, height, size=12, bold=False, color=DARK):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return txBox

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("document_title", "SOP Training Module")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Generated Training Module"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xA0, 0xC4, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Summary ─────────────────────────────────────────────────────
    summary = data.get("summary", {})
    slide = add_slide_with_header("📋 Summary")

    overview = summary.get("overview", "")
    add_text_box(slide, overview, 0.5, 1.4, 12, 1.2, size=13)

    # Objectives
    if summary.get("key_objectives"):
        add_text_box(slide, "Key Objectives", 0.5, 2.7, 5.5, 0.4, size=12, bold=True, color=ACCENT)
        objs = "\n".join(f"• {o}" for o in summary["key_objectives"][:5])
        add_text_box(slide, objs, 0.5, 3.1, 5.5, 3.5, size=11)

    # Rules
    if summary.get("important_rules"):
        add_text_box(slide, "Important Rules", 7.0, 2.7, 5.5, 0.4, size=12, bold=True, color=ACCENT)
        rules = "\n".join(f"• {r}" for r in summary["important_rules"][:5])
        add_text_box(slide, rules, 7.0, 3.1, 5.5, 3.5, size=11)

    # ── Training Step Slides ──────────────────────────────────────────────────
    for step in data.get("training_steps", []):
        num = step.get("step_number", "")
        step_title = step.get("title", "")
        slide = add_slide_with_header(f"Step {num}: {step_title}", ACCENT)

        add_text_box(slide, step.get("description", ""), 0.5, 1.4, 12, 2.5, size=13)

        if step.get("example"):
            add_text_box(slide, "💡 Example", 0.5, 4.0, 3, 0.4, size=11, bold=True, color=ACCENT)
            add_text_box(slide, step["example"], 0.5, 4.4, 12, 2, size=11)

    # ── Quiz Slides ───────────────────────────────────────────────────────────
    for q in data.get("quiz", []):
        qnum = q.get("question_number", "")
        qtype = q.get("type", "mcq").upper()
        slide = add_slide_with_header(f"Q{qnum} ({qtype})", RGBColor(0x06, 0x4E, 0x3B))

        add_text_box(slide, q.get("question", ""), 0.5, 1.4, 12, 1.5, size=14, bold=True)

        opts = q.get("options", [])
        opts_text = "\n".join(opts)
        add_text_box(slide, opts_text, 0.5, 3.0, 8, 2.5, size=12)

        answer_text = f"✅ Answer: {q.get('answer', '')} — {q.get('explanation', '')}"
        add_text_box(slide, answer_text, 0.5, 5.6, 12, 1.5, size=11, color=RGBColor(0x06, 0x4E, 0x3B))

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
